import os
import PyPDF2
import nltk
nltk.download('punkt', quiet=True)
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
from sumy.nlp.stemmers import Stemmer
from sumy.utils import get_stop_words

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import pyttsx3
from moviepy.editor import *
from PIL import Image, ImageDraw, ImageFont

def extract_text_from_pdf(pdf_path):
    with open(pdf_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        num_pages = len(pdf_reader.pages)
        all_text = ""
        
        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            text = page.extract_text()
            all_text += f"{text}\n\n"
    return all_text

def summarize_text(text, language="english"):
    parser = PlaintextParser.from_string(text, Tokenizer(language))
    stemmer = Stemmer(language)
    summarizer = LsaSummarizer(stemmer)
    summarizer.stop_words = get_stop_words(language)
    
    num_sentences = len(nltk.sent_tokenize(text))
    summary_sentences = max(int(num_sentences * 0.2), 5)
    
    summary = summarizer(parser.document, summary_sentences)
    return " ".join(str(sentence) for sentence in summary)

def create_slides_from_summary(title, content, output_folder):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    blank_slide_layout = prs.slide_layouts[6]

    # Greeting slide
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_content(prs, slide, f"Hello and welcome to the quick summary of {title}", is_title_slide=True)

    # Content slides
    words = content.split()
    chunks = []
    current_chunk = []
    
    for word in words:
        current_chunk.append(word)
        if len(' '.join(current_chunk)) > 200:
            chunks.append(' '.join(current_chunk))
            current_chunk = []
    if current_chunk:
        chunks.append(' '.join(current_chunk))

    for i, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(blank_slide_layout)
        create_slide_content(prs, slide, chunk, slide_number=i+1)

    # Closing slide
    slide = prs.slides.add_slide(blank_slide_layout)
    create_slide_content(prs, slide, "Thanks for watching and see you next time", is_title_slide=True)

    pptx_path = os.path.join(output_folder, 'summary_presentation.pptx')
    prs.save(pptx_path)
    return pptx_path

def create_slide_content(prs, slide, content, is_title_slide=False, slide_number=None):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background

    left = top = Inches(0.1)
    width = prs.slide_width - Inches(0.2)
    height = prs.slide_height - Inches(0.2)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.color.rgb = RGBColor(255, 255, 255)  # White border
    shape.line.width = Pt(3)
    shape.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), prs.slide_height - Inches(1))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = content
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.font.size = Pt(44) if is_title_slide else Pt(28)
    p.font.bold = True

    if not is_title_slide:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.5))
        title_p = title_box.text_frame.add_paragraph()
        title_p.text = f"Slide {slide_number}"
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.font.size = Pt(20)
        title_p.font.bold = True

        spacing_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), prs.slide_width - Inches(1), Inches(0.3))
        spacing_p = spacing_box.text_frame.add_paragraph()
        spacing_p.text = " "

def text_to_speech(text, output_file='output.mp3', rate=150, volume=1.0, voice=1):
    engine = pyttsx3.init()
    engine.setProperty('rate', rate)
    engine.setProperty('volume', volume)
    voices = engine.getProperty('voices')
    if voice < len(voices):
        engine.setProperty('voice', voices[voice].id)
    engine.save_to_file(text, output_file)
    engine.runAndWait()

def create_video_slide(text, slide_number=None, is_title_slide=False, size=(1280, 720)):
    img = Image.new('RGB', size, color=(0, 0, 0))  # Black background
    draw = ImageDraw.Draw(img)
    
    draw.rectangle([10, 10, size[0]-10, size[1]-10], outline=(255, 255, 255), width=3)  # White border

    font_title = ImageFont.truetype("arial.ttf", 44 if is_title_slide else 28)
    font_content = ImageFont.truetype("arial.ttf", 28)
    
    if is_title_slide:
        lines = text.split('\n')
        y_text = (size[1] - len(lines) * 50) // 2
        for line in lines:
            bbox = font_title.getbbox(line)
            text_width = bbox[2] - bbox[0]
            x_text = (size[0] - text_width) // 2
            draw.text((x_text, y_text), line, font=font_title, fill=(255, 255, 255))  # White text
            y_text += 50
    else:
        draw.text((size[0]//2, 30), f"Slide {slide_number}", font=font_title, fill=(255, 255, 255), anchor="mt")  # White text
        
        words = text.split()
        lines = []
        current_line = []
        for word in words:
            current_line.append(word)
            line = ' '.join(current_line)
            bbox = font_content.getbbox(line)
            if bbox[2] > size[0] - 100:
                lines.append(' '.join(current_line[:-1]))
                current_line = [word]
        if current_line:
            lines.append(' '.join(current_line))
        
        y_text = 100
        for line in lines:
            bbox = font_content.getbbox(line)
            text_width = bbox[2] - bbox[0]
            x_text = (size[0] - text_width) // 2
            draw.text((x_text, y_text), line, font=font_content, fill=(255, 255, 255))  # White text
            y_text += 40
    
    return img

def create_video(title, summary, audio_file, output_file):
    audio_clip = AudioFileClip(audio_file)
    duration = audio_clip.duration
    
    clips = []

    # Greeting slide
    greeting_img = create_video_slide(f"Hello and welcome to the quick summary of {title}", is_title_slide=True)
    greeting_path = "temp_greeting.png"
    greeting_img.save(greeting_path)
    greeting_clip = ImageClip(greeting_path).set_duration(5)
    clips.append(greeting_clip)

    # Content slides
    sentences = nltk.sent_tokenize(summary)
    total_chars = sum(len(s) for s in sentences)
    
    current_time = 5  # Start after the greeting slide
    for i, sentence in enumerate(sentences):
        slide_duration = (len(sentence) / total_chars) * (duration - 10)  # Subtract time for greeting and closing slides
        img = create_video_slide(sentence, i+1)
        
        slide_path = f"temp_slide_{i}.png"
        img.save(slide_path)
        
        clip = ImageClip(slide_path).set_duration(slide_duration)
        clip = clip.set_start(current_time)
        clips.append(clip)
        
        current_time += slide_duration
        
        os.remove(slide_path)

    # Closing slide
    closing_img = create_video_slide("Thanks for watching and see you next time", is_title_slide=True)
    closing_path = "temp_closing.png"
    closing_img.save(closing_path)
    closing_clip = ImageClip(closing_path).set_duration(5)
    closing_clip = closing_clip.set_start(current_time)
    clips.append(closing_clip)

    video = CompositeVideoClip(clips)
    video = video.set_audio(audio_clip)
    
    video = video.fadeout(1).fadein(1)
    
    video.write_videofile(output_file, fps=24)

    # Clean up temporary files
    os.remove(greeting_path)
    os.remove(closing_path)

def process_pdf(pdf_path, output_folder):
    extracted_text = extract_text_from_pdf(pdf_path)
    summary = summarize_text(extracted_text)

    pdf_filename = os.path.splitext(os.path.basename(pdf_path))[0]
    title = f"{pdf_filename}"
    pptx_path = create_slides_from_summary(title, summary, output_folder)

    full_text = f"Hello and welcome to the quick summary of {title}. {summary} Thanks for watching and see you next time"
    audio_path = os.path.join(output_folder, 'summary_audio.mp3')
    text_to_speech(full_text, audio_path)

    video_path = os.path.join(output_folder, 'summary_video.mp4')
    create_video(title, summary, audio_path, video_path)

    return video_path, pptx_path

def clean_up_files(upload_folder, output_folder):
    for file in os.listdir(upload_folder):
        os.remove(os.path.join(upload_folder, file))
    for file in os.listdir(output_folder):
        os.remove(os.path.join(output_folder, file))